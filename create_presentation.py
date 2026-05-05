"""
Generate TikTok Viral Detector — NLP Models & Tatiana Results PowerPoint
Run: python create_presentation.py
Output: TikTok_Viral_Detector_Presentation.pptx
"""

from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).parent))

# ── Colour palette ─────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x0D, 0x1A)   # very dark navy
PANEL    = RGBColor(0x16, 0x16, 0x2E)   # slightly lighter panel
RED      = RGBColor(0xFF, 0x6B, 0x6B)   # viral red
YELLOW   = RGBColor(0xFF, 0xD9, 0x3D)   # strong hook yellow
GREEN    = RGBColor(0x6B, 0xCB, 0x77)   # solid green
BLUE     = RGBColor(0x4D, 0x96, 0xFF)   # accent blue
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0x88, 0x88, 0x99)
DIMWHITE = RGBColor(0xCC, 0xCC, 0xDD)
BORDER   = RGBColor(0x33, 0x33, 0x55)

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# Low-level drawing helpers
# ══════════════════════════════════════════════════════════════════════════════

def set_bg(slide, color: RGBColor = BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=PANEL, alpha=None, line_color=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bar(slide, x, y, bar_w_max, bar_h, value, max_val,
            color=RED, label="", label_color=WHITE, show_pct=True):
    """Draw a single horizontal bar."""
    fill_w = max(int(bar_w_max * value / max(max_val, 0.001)), 2)
    # Background track
    add_rect(slide, x, y, bar_w_max, bar_h, fill=RGBColor(0x22,0x22,0x44))
    # Fill
    add_rect(slide, x, y, fill_w, bar_h, fill=color)
    # Label on left
    if label:
        add_text(slide, label, x - Inches(2.2), y - Pt(1), Inches(2.1), bar_h + Pt(4),
                 size=10, color=DIMWHITE, align=PP_ALIGN.RIGHT)
    # Percentage on right
    if show_pct:
        pct_str = f"{round(value*100)}%"
        add_text(slide, pct_str, x + fill_w + Pt(4), y - Pt(1), Inches(0.5), bar_h + Pt(4),
                 size=9, bold=True, color=color)


def pill(slide, text, x, y, w=Inches(1.4), h=Inches(0.28), bg=RED, txt_color=WHITE, size=11):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, text, x, y + Pt(2), w, h - Pt(2),
             size=size, bold=True, color=txt_color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Data
# ══════════════════════════════════════════════════════════════════════════════

FEAT1 = [
    ("repetition_score",  0.4623),
    ("is_repeated",       0.3918),
    ("position_ratio",    0.0981),
    ("word_count_norm",   0.0203),
    ("positive_density",  0.0077),
    ("brevity",           0.0063),
    ("energy_density",    0.0038),
    ("rhyme_next",        0.0031),
    ("rhyme_prev",        0.0031),
    ("negative_density",  0.0025),
    ("strong_start",      0.0008),
]

FEAT3 = [
    ("line_count_norm",      0.1419),
    ("hook_line_ratio",      0.1377),
    ("top_viral_prob",       0.1368),
    ("avg_negative",         0.0865),
    ("avg_energy",           0.0827),
    ("avg_brevity",          0.0762),
    ("avg_viral_prob",       0.0588),
    ("max_repeats_norm",     0.0545),
    ("viral_line_density",   0.0513),
    ("strong_hook_density",  0.0505),
    ("repetition_rate",      0.0501),
    ("avg_positive",         0.0446),
    ("chorus_repeat_ratio",  0.0283),
]

TATIANA = [
    ("Fresh Meat",          1.00, "🔥", "#ff6b6b", "Rotting here down in the belly of the beast", "[Verse]",    97),
    ("Make Me Feel",        0.99, "🔥", "#ff6b6b", "Caught here",                                 "[Chorus]",   82),
    ("Save Me",             0.99, "🔥", "#ff6b6b", "You'll be waiting for someone to save you",   "[Chorus]",   88),
    ("Born Again Sinner",   0.98, "🔥", "#ff6b6b", "I'm a born again sinner",                     "[Chorus]",   94),
    ("Without You",         0.93, "🔥", "#ff6b6b", "But see without you it means nothing",        "[Chorus]",   90),
    ("You Make Me",         0.88, "🔥", "#ff6b6b", "You, You Make Me",                            "[Chorus]",   96),
    ("Anxiety - Acoustic",  0.63, "⚡", "#ffd93d", "C'mon and try and shake it",                  "[Chorus]",   75),
    ("Too Much",            0.58, "⚡", "#ffd93d", "Too much it ties me up",                      "[Chorus]",   89),
    ("Gallery",             0.44, "📉", "#888899", "when you're afraid and all alone?",           "[Tag]",      71),
    ("What It Is About You",0.32, "📉", "#888899", "With your guilt free heart",                  "[Chorus]",   97),
    ("AWS",                 0.17, "📉", "#888899", "I finally found my way into my mind",         "[Chorus]",   78),
    ("London Don't Lie",    0.09, "📉", "#888899", "No don't lie no no",                          "[Chorus]",   68),
]


# ══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    # Big gradient-ish accent bar on left
    add_rect(slide, 0, 0, Inches(0.18), H, fill=RED)

    # TikTok logo-ish double accent
    add_rect(slide, Inches(0.18), 0, Inches(0.06), H, fill=BLUE)

    # Title
    add_text(slide, "TikTok Viral Detector", Inches(0.55), Inches(1.8),
             Inches(8), Inches(1.2), size=44, bold=True, color=WHITE)

    # Sub
    add_text(slide,
             "NLP Model Architecture · Training Data · Feature Weights\n"
             "Tatiana DeMaria — Full Catalog Viral Analysis",
             Inches(0.55), Inches(3.1), Inches(9), Inches(1.0),
             size=20, color=DIMWHITE)

    # Pills
    for i, (txt, col) in enumerate([("Model 1: Line Classifier", RED),
                                     ("Model 2: Similarity", BLUE),
                                     ("Model 3: Song Virality", YELLOW)]):
        pill(slide, txt, Inches(0.55 + i * 3.1), Inches(5.5),
             w=Inches(2.85), h=Inches(0.38), bg=col,
             txt_color=RGBColor(0,0,0) if col==YELLOW else WHITE, size=12)

    # Right decoration — faint lines
    for i in range(8):
        w_line = Inches(3.5 - i * 0.35)
        add_rect(slide, W - w_line, Inches(0.6 + i * 0.75),
                 w_line, Inches(0.06),
                 fill=RGBColor(0x33, 0x33, 0x66))


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text(slide, "How the System Works", Inches(0.5), Inches(0.25),
             Inches(10), Inches(0.65), size=30, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(5.5), Inches(0.04), fill=RED)

    steps = [
        (RED,    "INPUT",         "User pastes any song lyrics\n(or Tatiana's CSV is loaded directly)"),
        (BLUE,   "PARSE",         "Section headers detected & stripped\n(Chorus, Verse, Bridge, etc.)"),
        (YELLOW, "MODEL 1",       "Every lyric line scored 0–100%\nGradient Boosting on 11 features"),
        (GREEN,  "MODEL 3",       "Full song scored 0–100%\n13 aggregated song-level features"),
        (RED,    "OUTPUT",        "Viral line + TikTok clip\nSong score + ranked table"),
    ]

    box_w = Inches(2.3)
    box_h = Inches(1.6)
    gap   = Inches(0.15)
    start_x = Inches(0.3)
    y = Inches(1.6)

    for i, (color, title, body) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_rect(slide, x, y, box_w, box_h, fill=PANEL, line_color=color, line_w=Pt(2))
        add_text(slide, title, x, y + Inches(0.1), box_w, Inches(0.4),
                 size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + Inches(0.08), y + Inches(0.48),
                 box_w - Inches(0.16), Inches(1.0),
                 size=10, color=DIMWHITE, align=PP_ALIGN.CENTER)
        # Arrow
        if i < len(steps) - 1:
            ax = x + box_w + Pt(4)
            add_text(slide, "▶", ax, y + Inches(0.6), Inches(0.18), Inches(0.5),
                     size=16, bold=True, color=GREY, align=PP_ALIGN.CENTER)

    # Training data section
    add_text(slide, "Training Data", Inches(0.5), Inches(3.55),
             Inches(8), Inches(0.45), size=18, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(3.95), Inches(12.0), Inches(0.04), fill=BORDER)

    data_items = [
        (BLUE,   "Summer Dataset",    "4,785 songs  ·  5 tiers (0–4)\nTier 4 = Super Hits"),
        (GREEN,  "Balanced Catalog",  "50 songs/tier sampled\n= 250 balanced songs"),
        (YELLOW, "Genius Lyrics API", "Lyrics fetched for 231 songs\nLyric text = model input"),
        (RED,    "TikTok Trending",   "6 live trending songs\nLabelled viral = 1"),
    ]

    bw  = Inches(2.85)
    for i, (color, title, body) in enumerate(data_items):
        x = Inches(0.35) + i * (bw + Inches(0.2))
        add_rect(slide, x, Inches(4.1), bw, Inches(1.6), fill=PANEL, line_color=color, line_w=Pt(1.5))
        add_text(slide, title, x + Inches(0.1), Inches(4.15), bw - Inches(0.2), Inches(0.4),
                 size=12, bold=True, color=color)
        add_text(slide, body, x + Inches(0.1), Inches(4.55), bw - Inches(0.2), Inches(1.0),
                 size=10, color=DIMWHITE)


def slide_model1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text(slide, "Model 1 — Viral Line Classifier", Inches(0.5), Inches(0.25),
             Inches(10), Inches(0.65), size=30, bold=True, color=RED)
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(5.5), Inches(0.04), fill=RED)

    # Left: specs
    specs = [
        ("Algorithm",    "Gradient Boosting Classifier"),
        ("n_estimators", "300 trees"),
        ("max_depth",    "4"),
        ("learning_rate","0.08"),
        ("subsample",    "0.85"),
        ("Training data","12,282 lyric lines"),
        ("Classes",      "Viral (1) / Non-Viral (0)"),
        ("CV F1 Score",  "0.845 ± 0.008"),
        ("Accuracy",     "87.3%"),
    ]

    add_text(slide, "Hyperparameters & Results", Inches(0.5), Inches(1.1),
             Inches(5), Inches(0.35), size=14, bold=True, color=WHITE)

    for i, (k, v) in enumerate(specs):
        y = Inches(1.55) + i * Inches(0.52)
        add_rect(slide, Inches(0.5), y, Inches(5.3), Inches(0.46),
                 fill=PANEL if i % 2 == 0 else RGBColor(0x12,0x12,0x28))
        add_text(slide, k, Inches(0.6), y + Pt(4), Inches(1.8), Inches(0.38),
                 size=11, bold=True, color=GREY)
        add_text(slide, v, Inches(2.45), y + Pt(4), Inches(3.2), Inches(0.38),
                 size=11, color=WHITE)

    # Labelling rule
    add_rect(slide, Inches(0.5), Inches(6.4), Inches(5.3), Inches(0.72),
             fill=PANEL, line_color=YELLOW, line_w=Pt(1.5))
    add_text(slide, "Weak Supervision Label Rule:",
             Inches(0.6), Inches(6.45), Inches(5.0), Inches(0.28),
             size=11, bold=True, color=YELLOW)
    add_text(slide, "Line is VIRAL (y=1) if it is in [Chorus/Hook/Bridge] section  "
             "OR  repeats ≥ 3× in the song",
             Inches(0.6), Inches(6.72), Inches(5.0), Inches(0.32),
             size=10, color=DIMWHITE)

    # Right: class breakdown
    add_text(slide, "Training Set Breakdown", Inches(6.1), Inches(1.1),
             Inches(6.5), Inches(0.35), size=14, bold=True, color=WHITE)

    add_rect(slide, Inches(6.1), Inches(1.5), Inches(6.5), Inches(1.35), fill=PANEL)
    add_text(slide, "12,282 lyric lines total", Inches(6.2), Inches(1.55),
             Inches(6.0), Inches(0.35), size=12, bold=True, color=WHITE)
    add_text(slide, "49%  Viral lines  (6,000)",
             Inches(6.2), Inches(1.9), Inches(3.0), Inches(0.35), size=11, color=RED, bold=True)
    add_text(slide, "51%  Non-viral lines  (6,282)",
             Inches(9.4), Inches(1.9), Inches(3.0), Inches(0.35), size=11, color=GREY)

    # Feature importance chart title
    add_text(slide, "Feature Importance Weights (Model 1)", Inches(6.1), Inches(3.0),
             Inches(6.8), Inches(0.4), size=14, bold=True, color=WHITE)

    bar_x     = Inches(8.5)
    bar_w_max = Inches(4.3)
    bar_h     = Inches(0.30)
    gap       = Inches(0.09)

    for i, (name, imp) in enumerate(FEAT1):
        y = Inches(3.5) + i * (bar_h + gap)
        col = RED if imp >= 0.1 else (YELLOW if imp >= 0.03 else BLUE)
        add_bar(slide, bar_x, y, bar_w_max, bar_h, imp, FEAT1[0][1],
                color=col, label=name, show_pct=True)


def slide_model2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text(slide, "Model 2 — Song Similarity Engine", Inches(0.5), Inches(0.25),
             Inches(10), Inches(0.65), size=30, bold=True, color=BLUE)
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(5.5), Inches(0.04), fill=BLUE)

    steps = [
        (BLUE,   "TF-IDF\nVectorizer",    "Converts each song's\nlyrics into a numeric\nvector (25,000 terms)\nbigrams (1,2)\nsublinear_tf = True"),
        (GREEN,  "Cosine\nSimilarity",    "NearestNeighbors\nwith cosine metric\n(brute-force, fast\nfor 231 songs)"),
        (YELLOW, "Token\nOverlap",        "Raw unigram sets\nper song computed\nat train time\n(bigrams too specific)"),
        (RED,    "Why-Match\nExplanation","Shared vocabulary\ndisplayed as keyword\ntags in the UI"),
    ]

    bw = Inches(2.8)
    for i, (color, title, body) in enumerate(steps):
        x = Inches(0.4) + i * (bw + Inches(0.25))
        add_rect(slide, x, Inches(1.2), bw, Inches(2.4), fill=PANEL, line_color=color, line_w=Pt(2))
        add_text(slide, title, x, Inches(1.25), bw, Inches(0.65),
                 size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + Inches(0.1), Inches(1.9), bw - Inches(0.2), Inches(1.6),
                 size=10, color=DIMWHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(slide, "→", x + bw + Inches(0.02), Inches(2.0), Inches(0.28), Inches(0.5),
                     size=18, bold=True, color=GREY, align=PP_ALIGN.CENTER)

    # Specs
    add_text(slide, "Vectorizer Settings", Inches(0.5), Inches(3.85),
             Inches(5), Inches(0.38), size=14, bold=True, color=WHITE)

    specs2 = [
        ("ngram_range",   "(1, 2)  — unigrams + bigrams"),
        ("max_features",  "25,000 vocabulary terms"),
        ("min_df",        "1  — keep all terms"),
        ("max_df",        "0.90  — exclude near-universal terms"),
        ("sublinear_tf",  "True  — log-scale term frequencies"),
        ("strip_accents", "unicode — normalise foreign characters"),
        ("Catalog size",  "231 songs with lyrics"),
    ]

    for i, (k, v) in enumerate(specs2):
        y = Inches(4.3) + i * Inches(0.43)
        add_rect(slide, Inches(0.5), y, Inches(5.8), Inches(0.38),
                 fill=PANEL if i % 2 == 0 else RGBColor(0x12,0x12,0x28))
        add_text(slide, k, Inches(0.6), y + Pt(3), Inches(1.6), Inches(0.32),
                 size=10, bold=True, color=GREY)
        add_text(slide, v, Inches(2.25), y + Pt(3), Inches(3.9), Inches(0.32),
                 size=10, color=WHITE)

    # Note box
    add_rect(slide, Inches(6.5), Inches(3.85), Inches(6.3), Inches(3.2),
             fill=PANEL, line_color=BLUE, line_w=Pt(1.5))
    add_text(slide, "How Similarity is Used in the App",
             Inches(6.6), Inches(3.95), Inches(6.0), Inches(0.35),
             size=13, bold=True, color=BLUE)
    note_body = (
        "1.  Paste lyrics into Tab 1\n"
        "2.  TF-IDF keyword extraction identifies\n"
        "     the most distinctive words in the lyrics\n\n"
        "3.  Those keywords are sent to the TikTok\n"
        "     RapidAPI to find currently trending songs\n\n"
        "4.  Results ranked by trend_score:\n"
        "     0.80 × log(play_count)\n"
        "   + 0.20 × keyword_match_ratio\n\n"
        "5.  Each match shows WHY it matched\n"
        "     (shared keywords or lyric search terms)"
    )
    add_text(slide, note_body, Inches(6.6), Inches(4.4), Inches(6.0), Inches(2.5),
             size=10, color=DIMWHITE)


def slide_model3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text(slide, "Model 3 — Song-Level TikTok Virality Predictor",
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.65),
             size=28, bold=True, color=YELLOW)
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(6.0), Inches(0.04), fill=YELLOW)

    # Left: specs
    specs3 = [
        ("Algorithm",     "Gradient Boosting Classifier"),
        ("n_estimators",  "200 trees"),
        ("max_depth",     "3"),
        ("learning_rate", "0.08"),
        ("subsample",     "0.80"),
        ("Training songs","187 songs (tiers 0,1,3,4,-1)"),
        ("Label 1 (viral)","Tier 3/4 + TikTok trending (tier -1)"),
        ("Label 0 (not)",  "Tier 0/1 (non-hits / minor hits)"),
        ("Tier 2",         "Skipped (ambiguous mid-tier)"),
        ("CV F1 Score",   "0.678 ± 0.057"),
    ]

    add_text(slide, "Hyperparameters & Label Design", Inches(0.5), Inches(1.1),
             Inches(5.5), Inches(0.35), size=14, bold=True, color=WHITE)

    for i, (k, v) in enumerate(specs3):
        y = Inches(1.55) + i * Inches(0.48)
        add_rect(slide, Inches(0.5), y, Inches(5.6), Inches(0.43),
                 fill=PANEL if i % 2 == 0 else RGBColor(0x12,0x12,0x28))
        add_text(slide, k, Inches(0.6), y + Pt(3), Inches(2.1), Inches(0.37),
                 size=10, bold=True, color=GREY)
        add_text(slide, v, Inches(2.75), y + Pt(3), Inches(3.2), Inches(0.37),
                 size=10, color=WHITE)

    # Right: feature importance chart
    add_text(slide, "Feature Importance Weights (Model 3)", Inches(6.3), Inches(1.1),
             Inches(6.8), Inches(0.4), size=14, bold=True, color=WHITE)

    bar_x     = Inches(9.1)
    bar_w_max = Inches(3.9)
    bar_h     = Inches(0.27)
    gap       = Inches(0.07)

    for i, (name, imp) in enumerate(FEAT3):
        y = Inches(1.6) + i * (bar_h + gap)
        col = YELLOW if imp >= 0.12 else (GREEN if imp >= 0.07 else BLUE)
        add_bar(slide, bar_x, y, bar_w_max, bar_h, imp, FEAT3[0][1],
                color=col, label=name, show_pct=True)

    # How it works note
    add_rect(slide, Inches(6.3), Inches(5.9), Inches(6.7), Inches(1.35),
             fill=PANEL, line_color=YELLOW, line_w=Pt(1.5))
    add_text(slide, "Dependency Chain:",
             Inches(6.4), Inches(5.95), Inches(6.3), Inches(0.3),
             size=11, bold=True, color=YELLOW)
    add_text(slide,
             "Model 1 must be trained first.  Model 3 calls Model 1 on every training "
             "song to compute avg_viral_prob, top_viral_prob, and density features — "
             "these song-level aggregates are then fed into Model 3.",
             Inches(6.4), Inches(6.28), Inches(6.3), Inches(0.85),
             size=10, color=DIMWHITE)


def slide_tatiana_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text(slide, "Tatiana DeMaria — Catalog Overview",
             Inches(0.5), Inches(0.25), Inches(10), Inches(0.65),
             size=30, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(5.5), Inches(0.04), fill=RED)

    # Big stat boxes
    stats = [
        (RED,    "6",  "🔥 High Potential\n(≥65%)"),
        (YELLOW, "2",  "⚡ Moderate\n(45–64%)"),
        (GREY,   "4",  "📉 Low\n(<45%)"),
        (GREEN,  "58%","Avg Viral Score"),
        (BLUE,   "12", "Songs Analysed"),
    ]

    bw = Inches(2.3)
    for i, (color, num, label) in enumerate(stats):
        x = Inches(0.3) + i * (bw + Inches(0.17))
        add_rect(slide, x, Inches(1.2), bw, Inches(1.55), fill=PANEL, line_color=color, line_w=Pt(2))
        add_text(slide, num, x, Inches(1.3), bw, Inches(0.8),
                 size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, Inches(2.1), bw, Inches(0.55),
                 size=11, color=DIMWHITE, align=PP_ALIGN.CENTER)

    # Top pick highlight
    add_text(slide, "Top Pick",
             Inches(0.5), Inches(3.0), Inches(3), Inches(0.4),
             size=16, bold=True, color=RED)
    add_rect(slide, Inches(0.5), Inches(3.45), Inches(5.7), Inches(1.5),
             fill=PANEL, line_color=RED, line_w=Pt(2))
    add_text(slide, "Fresh Meat", Inches(0.6), Inches(3.5), Inches(5.5), Inches(0.55),
             size=22, bold=True, color=RED)
    add_text(slide, "TikTok Viral Score: 100%  ·  🔥 High TikTok Potential",
             Inches(0.6), Inches(4.05), Inches(5.5), Inches(0.35),
             size=11, color=DIMWHITE)
    add_text(slide, '"Rotting here down in the belly of the beast"',
             Inches(0.6), Inches(4.4), Inches(5.5), Inches(0.35),
             size=11, italic=True, color=YELLOW)

    # Distribution bar
    add_text(slide, "Score Distribution", Inches(6.5), Inches(3.0),
             Inches(6.0), Inches(0.4), size=16, bold=True, color=WHITE)

    dist_items = [
        ("Fresh Meat",          1.00, RED),
        ("Make Me Feel",        0.99, RED),
        ("Save Me",             0.99, RED),
        ("Born Again Sinner",   0.98, RED),
        ("Without You",         0.93, RED),
        ("You Make Me",         0.88, RED),
        ("Anxiety - Acoustic",  0.63, YELLOW),
        ("Too Much",            0.58, YELLOW),
        ("Gallery",             0.44, GREY),
        ("What It Is About You",0.32, GREY),
        ("AWS",                 0.17, GREY),
        ("London Don't Lie",    0.09, GREY),
    ]

    bar_x     = Inches(9.0)
    bar_w_max = Inches(3.8)
    bar_h     = Inches(0.28)
    gap       = Inches(0.05)

    for i, (name, prob, color) in enumerate(dist_items):
        y = Inches(3.5) + i * (bar_h + gap)
        add_bar(slide, bar_x, y, bar_w_max, bar_h, prob, 1.0,
                color=color, label=name, show_pct=True)


def slide_tatiana_ranked(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text(slide, "Tatiana DeMaria — Full Ranked Table",
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(0.82), Inches(12.5), Inches(0.04), fill=RED)

    # Column headers
    headers = ["#", "Song", "Score", "Rating", "Viral Line", "Section", "Line %"]
    col_x   = [Inches(0.35), Inches(0.75), Inches(3.65), Inches(4.55),
               Inches(5.65), Inches(10.5), Inches(11.65)]
    col_w   = [Inches(0.38), Inches(2.8), Inches(0.85), Inches(0.95),
               Inches(4.75), Inches(1.1), Inches(0.9)]

    hdr_y = Inches(1.0)
    add_rect(slide, Inches(0.35), hdr_y, Inches(12.6), Inches(0.38), fill=BORDER)
    for hdr, x, w in zip(headers, col_x, col_w):
        add_text(slide, hdr, x, hdr_y + Pt(2), w, Inches(0.32),
                 size=10, bold=True, color=GREY, align=PP_ALIGN.LEFT)

    row_h = Inches(0.47)
    for i, (song, prob, emoji, hex_color, viral_line, section, line_pct) in enumerate(TATIANA):
        y = Inches(1.42) + i * row_h
        row_bg = PANEL if i % 2 == 0 else RGBColor(0x12, 0x12, 0x28)
        add_rect(slide, Inches(0.35), y, Inches(12.6), row_h - Pt(2), fill=row_bg)

        # Rank
        add_text(slide, str(i+1), col_x[0], y + Pt(4), col_w[0], row_h,
                 size=10, color=GREY, align=PP_ALIGN.CENTER)
        # Song name
        add_text(slide, song, col_x[1], y + Pt(4), col_w[1], row_h,
                 size=10, bold=True, color=WHITE)
        # Score
        color_obj = RGBColor.from_string(hex_color[1:])
        add_text(slide, f"{round(prob*100)}%", col_x[2], y + Pt(4), col_w[2], row_h,
                 size=12, bold=True, color=color_obj, align=PP_ALIGN.CENTER)
        # Rating emoji
        add_text(slide, emoji, col_x[3], y + Pt(4), col_w[3], row_h,
                 size=14, align=PP_ALIGN.CENTER)
        # Viral line (truncate)
        vl_trunc = viral_line[:62] + "…" if len(viral_line) > 62 else viral_line
        add_text(slide, f'"{vl_trunc}"', col_x[4], y + Pt(4), col_w[4], row_h,
                 size=9, italic=True, color=DIMWHITE)
        # Section
        add_text(slide, section, col_x[5], y + Pt(4), col_w[5], row_h,
                 size=9, color=GREY, align=PP_ALIGN.CENTER)
        # Line pct
        add_text(slide, f"{line_pct}%", col_x[6], y + Pt(4), col_w[6], row_h,
                 size=10, bold=True, color=color_obj, align=PP_ALIGN.CENTER)

    # Legend
    for txt, col in [("🔥 High ≥65%", RED), ("⚡ Moderate 45–64%", YELLOW), ("📉 Low <45%", GREY)]:
        pass  # done via emoji column

    add_rect(slide, Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.3),
             fill=PANEL)
    add_text(slide,
             "Score = Model 3 song-level viral probability  ·  "
             "Line % = Model 1 viral probability for the top clip line  ·  "
             "Section = where the viral line lives",
             Inches(0.45), Inches(7.12), Inches(12.3), Inches(0.26),
             size=8, color=GREY)


def slide_top_songs(prs):
    """Spotlight the top 3 songs with their viral lines and why."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text(slide, "Top 3 Songs — Viral Line Spotlight",
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.6),
             size=30, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(0.88), Inches(6.0), Inches(0.04), fill=RED)

    top3 = TATIANA[:3]
    box_w = Inches(4.1)
    gap   = Inches(0.2)

    for i, (song, prob, emoji, hex_color, viral_line, section, line_pct) in enumerate(top3):
        x     = Inches(0.3) + i * (box_w + gap)
        color = RGBColor.from_string(hex_color[1:])

        # Card
        add_rect(slide, x, Inches(1.1), box_w, Inches(6.1),
                 fill=PANEL, line_color=color, line_w=Pt(2.5))

        # Rank badge
        add_rect(slide, x, Inches(1.1), Inches(0.5), Inches(0.45), fill=color)
        add_text(slide, f"#{i+1}", x, Inches(1.1), Inches(0.5), Inches(0.45),
                 size=14, bold=True, color=WHITE if i > 0 else RGBColor(0,0,0),
                 align=PP_ALIGN.CENTER)

        # Song name
        add_text(slide, song, x + Inches(0.55), Inches(1.13),
                 box_w - Inches(0.6), Inches(0.55),
                 size=16, bold=True, color=color)

        # Score
        add_text(slide, f"{emoji} {round(prob*100)}%",
                 x + Inches(0.1), Inches(1.75), box_w - Inches(0.2), Inches(0.4),
                 size=13, bold=True, color=color)

        # Viral line box
        add_rect(slide, x + Inches(0.1), Inches(2.3), box_w - Inches(0.2), Inches(1.15),
                 fill=RGBColor(0x22,0x11,0x11))
        add_text(slide, "VIRAL LINE TO FILM",
                 x + Inches(0.18), Inches(2.35), box_w - Inches(0.3), Inches(0.28),
                 size=8, bold=True, color=RED)
        add_text(slide, f'"{viral_line}"',
                 x + Inches(0.18), Inches(2.63), box_w - Inches(0.3), Inches(0.75),
                 size=11, italic=True, color=WHITE)

        # Stats
        stat_items = [
            (f"Line Viral Prob", f"{line_pct}%"),
            (f"Best Section",    section),
            ("Song Score",       f"{round(prob*100)}%"),
        ]
        for j, (k, v) in enumerate(stat_items):
            sy = Inches(3.6) + j * Inches(0.55)
            add_rect(slide, x + Inches(0.1), sy, box_w - Inches(0.2), Inches(0.46),
                     fill=RGBColor(0x12,0x12,0x28))
            add_text(slide, k, x + Inches(0.18), sy + Pt(4),
                     Inches(1.5), Inches(0.35), size=10, color=GREY)
            add_text(slide, v, x + Inches(1.75), sy + Pt(4),
                     Inches(1.8), Inches(0.35), size=11, bold=True, color=color)

        # Why viral
        add_text(slide, "Why it will go viral:",
                 x + Inches(0.1), Inches(5.35), box_w - Inches(0.2), Inches(0.3),
                 size=10, bold=True, color=YELLOW)

        why_map = {
            "Fresh Meat":        "High energy language + strong repetition rate. "
                                 "The line \"belly of the beast\" is visceral and clip-worthy.",
            "Make Me Feel":      "Ultra-short viral line (2 words). Appears in Chorus. "
                                 "POV / lip sync format is obvious.",
            "Save Me":           "Emotionally resonant hook. 88% viral probability — "
                                 "strong repetition drives duet culture.",
        }
        add_text(slide, why_map.get(song, "Strong hook structure and repeating chorus."),
                 x + Inches(0.1), Inches(5.68), box_w - Inches(0.2), Inches(1.3),
                 size=9, color=DIMWHITE)


def slide_key_findings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text(slide, "Key Findings & Recommendations",
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.65),
             size=30, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(6.0), Inches(0.04), fill=GREEN)

    findings = [
        (RED,    "🔥", "6 / 12 songs have High TikTok Potential",
                 "Fresh Meat, Make Me Feel, Save Me, Born Again Sinner, "
                 "Without You, and You Make Me all score ≥88%."),
        (YELLOW, "⚡", "Repetition is the strongest signal",
                 "Model 1 gives 46% weight to repetition_score and 39% to is_repeated. "
                 "Songs where a hook line repeats ≥3× are consistently flagged viral."),
        (GREEN,  "✓", "Chorus lines dominate across all songs",
                 "10 of 12 top viral lines live in the Chorus section. "
                 "Prioritise chorus polish for maximum TikTok impact."),
        (BLUE,   "↑", "Short lines outperform long lines",
                 "\"Caught here\" (2 words, 99% song score) shows brevity wins. "
                 "Aim for 4–9 words in the hook line."),
        (GREY,   "▼", "London Don't Lie scores lowest (9%)",
                 "Low repetition rate, complex structure, limited hook density. "
                 "Consider a stripped-back acoustic version for TikTok."),
    ]

    for i, (color, icon, headline, detail) in enumerate(findings):
        y = Inches(1.15) + i * Inches(1.15)
        # Icon circle
        add_rect(slide, Inches(0.4), y + Inches(0.12), Inches(0.52), Inches(0.52),
                 fill=color)
        add_text(slide, icon, Inches(0.4), y + Inches(0.12), Inches(0.52), Inches(0.52),
                 size=18, align=PP_ALIGN.CENTER)
        # Headline
        add_text(slide, headline, Inches(1.1), y + Inches(0.05),
                 Inches(11.5), Inches(0.38), size=14, bold=True, color=color)
        # Detail
        add_text(slide, detail, Inches(1.1), y + Inches(0.44),
                 Inches(11.5), Inches(0.55), size=11, color=DIMWHITE)
        # Divider
        if i < len(findings) - 1:
            add_rect(slide, Inches(0.4), y + Inches(1.05), Inches(12.5), Inches(0.02),
                     fill=BORDER)

    # Footer
    add_rect(slide, 0, H - Inches(0.45), W, Inches(0.45), fill=RGBColor(0x11,0x11,0x22))
    add_text(slide, "TikTok Viral Detector  ·  Models: Gradient Boosting + TF-IDF  ·  "
             "Training: 231 catalog songs  ·  Tatiana: 12 songs analysed",
             Inches(0.4), H - Inches(0.4), W - Inches(0.8), Inches(0.35),
             size=8, color=GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Build & save
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Building slides …")
    slide_title(prs);            print("  1/8  Title")
    slide_pipeline(prs);         print("  2/8  Pipeline & Training Data")
    slide_model1(prs);           print("  3/8  Model 1 — Viral Line Classifier")
    slide_model2(prs);           print("  4/8  Model 2 — Similarity Engine")
    slide_model3(prs);           print("  5/8  Model 3 — Song Virality")
    slide_tatiana_overview(prs); print("  6/8  Tatiana Overview")
    slide_tatiana_ranked(prs);   print("  7/8  Full Ranked Table")
    slide_top_songs(prs);        print("  8/8  Top 3 Spotlight")
    slide_key_findings(prs);     print("  9/9  Key Findings")

    out = Path(__file__).parent.parent / "TikTok_Viral_Detector_Presentation.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")
    return out


if __name__ == "__main__":
    build()
